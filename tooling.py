from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Callable
import json
import os
import re
import subprocess
import urllib.error
import urllib.parse
import urllib.request
import zipfile
import xml.etree.ElementTree as ET


class ToolExecutionError(Exception):
    pass


def safe_resolve_path(root: Path, user_path: str) -> Path:
    candidate = (root / user_path).resolve() if not Path(user_path).is_absolute() else Path(user_path).resolve()
    try:
        candidate.relative_to(root)
    except ValueError as exc:
        raise ToolExecutionError(f"Path out of workspace: {candidate}; allowed root: {root}") from exc
    return candidate


@dataclass
class Tool:
    name: str
    description: str
    input_schema: dict[str, Any]
    handler: Callable[[dict[str, Any]], str]


class ToolRegistry:
    def __init__(self, root: Path):
        self.root = root
        self._tools: dict[str, Tool] = {}
        self._register_builtin_tools()

    def _register_builtin_tools(self) -> None:
        self.register(
            Tool(
                name="get_current_time",
                description="Get local time, or time at a specific UTC offset.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "utc_offset": {
                            "type": "string",
                            "description": "Optional UTC offset like +08:00 or -05:00.",
                        }
                    },
                    "additionalProperties": False,
                },
                handler=self._get_current_time,
            )
        )

        self.register(
            Tool(
                name="get_weather",
                description="Get current weather and short forecast for a location.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "location": {
                            "type": "string",
                            "description": "City or place name.",
                        }
                    },
                    "required": ["location"],
                    "additionalProperties": False,
                },
                handler=self._get_weather,
            )
        )

        self.register(
            Tool(
                name="search_web",
                description="Search the web for a topic or keyword and return concise results.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "Keyword, noun, or question.",
                        },
                        "max_results": {
                            "type": "integer",
                            "description": "Maximum number of results, default 5.",
                            "minimum": 1,
                            "maximum": 10,
                        },
                    },
                    "required": ["query"],
                    "additionalProperties": False,
                },
                handler=self._search_web,
            )
        )

        self.register(
            Tool(
                name="list_files",
                description="List files and sub-directories under a workspace-relative directory.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Workspace-relative directory path, default '.'.",
                        }
                    },
                    "required": [],
                    "additionalProperties": False,
                },
                handler=self._list_files,
            )
        )

        self.register(
            Tool(
                name="read_text_file",
                description="Read content from text-like files. Supports txt/md/json/csv/tsv/pdf/docx/xlsx/pptx.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Workspace-relative file path."},
                        "max_chars": {
                            "type": "integer",
                            "description": "Maximum chars to return, default 4000.",
                            "minimum": 100,
                            "maximum": 20000,
                        },
                    },
                    "required": ["path"],
                    "additionalProperties": False,
                },
                handler=self._read_text_file,
            )
        )

        self.register(
            Tool(
                name="write_text_file",
                description="Write text content to a workspace file, overwrite or append.",
                input_schema={
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Workspace-relative file path."},
                        "content": {"type": "string", "description": "Text content to write."},
                        "mode": {
                            "type": "string",
                            "enum": ["overwrite", "append"],
                            "description": "overwrite replaces, append appends.",
                        },
                    },
                    "required": ["path", "content"],
                    "additionalProperties": False,
                },
                handler=self._write_text_file,
            )
        )
        self.register(
            Tool(
                name="run_shell_command",
                description=(
                    "Run a strictly validated shell command from an allow-list. "
                    "Only safe read-oriented commands are supported."
                ),
                input_schema={
                    "type": "object",
                    "properties": {
                        "command": {
                            "type": "array",
                            "description": "Command and args as tokenized string list, e.g. ['python','--version'].",
                            "items": {"type": "string"},
                            "minItems": 1,
                            "maxItems": 20,
                        },
                        "timeout_seconds": {
                            "type": "integer",
                            "minimum": 1,
                            "maximum": 20,
                            "description": "Execution timeout in seconds. Default 8.",
                        },
                    },
                    "required": ["command"],
                    "additionalProperties": False,
                },
                handler=self._run_shell_command,
            )
        )
        self.register(
            Tool(
                name="delegate_subagent",
                description=(
                    "Delegate a focused analysis subtask to a lightweight sub-agent. "
                    "Useful for summarizing uploaded files or extracting key points."
                ),
                input_schema={
                    "type": "object",
                    "properties": {
                        "task": {"type": "string", "description": "The subtask to delegate."},
                        "files": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Optional workspace-relative file paths for sub-agent analysis.",
                        },
                    },
                    "required": ["task"],
                    "additionalProperties": False,
                },
                handler=self._delegate_subagent,
            )
        )

    def register(self, tool: Tool) -> None:
        self._tools[tool.name] = tool

    def all_for_openai(self) -> list[dict[str, Any]]:
        return [
            {
                "type": "function",
                "function": {
                    "name": tool.name,
                    "description": tool.description,
                    "parameters": tool.input_schema,
                },
            }
            for tool in self._tools.values()
        ]

    def execute(self, name: str, raw_arguments: str) -> str:
        if name not in self._tools:
            raise ToolExecutionError(f"Unknown tool: {name}")

        try:
            arguments = json.loads(raw_arguments or "{}")
            if not isinstance(arguments, dict):
                raise ToolExecutionError("Tool arguments must be a JSON object")
        except json.JSONDecodeError as exc:
            raise ToolExecutionError(f"Invalid JSON arguments: {exc}") from exc

        return self._tools[name].handler(arguments)

    def _list_files(self, args: dict[str, Any]) -> str:
        rel_path = args.get("path", ".")
        path = safe_resolve_path(self.root, rel_path)
        if not path.exists():
            raise ToolExecutionError(f"Path does not exist: {rel_path}")
        if not path.is_dir():
            raise ToolExecutionError(f"Target is not a directory: {rel_path}")

        items = []
        for p in sorted(path.iterdir(), key=lambda x: x.name.lower()):
            kind = "dir" if p.is_dir() else "file"
            items.append({"name": p.name, "type": kind})
        return json.dumps({"path": str(path), "items": items}, ensure_ascii=False)

    def _read_text_file(self, args: dict[str, Any]) -> str:
        rel_path = args["path"]
        max_chars = int(args.get("max_chars", 4000))
        path = safe_resolve_path(self.root, rel_path)

        if not path.exists():
            raise ToolExecutionError(f"File does not exist: {rel_path}")
        if not path.is_file():
            raise ToolExecutionError(f"Target is not a file: {rel_path}")

        content, detected_format = self._read_file_content(path)
        clipped = content[:max_chars]
        return json.dumps(
            {
                "path": str(path),
                "detected_format": detected_format,
                "content": clipped,
                "truncated": len(content) > len(clipped),
                "total_chars": len(content),
            },
            ensure_ascii=False,
        )

    def _read_file_content(self, path: Path) -> tuple[str, str]:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self._read_pdf(path), "pdf"
        if suffix == ".docx":
            return self._read_docx(path), "docx"
        if suffix == ".xlsx":
            return self._read_xlsx(path), "xlsx"
        if suffix == ".pptx":
            return self._read_pptx(path), "pptx"
        if suffix in {".csv", ".tsv"}:
            return self._read_csv_like(path), suffix.lstrip(".")
        if suffix == ".json":
            return self._read_json(path), "json"

        raw = path.read_bytes()
        if self._looks_binary(raw):
            raise ToolExecutionError(
                f"Unsupported binary file for direct reading: {path.name}. Supported: text/JSON/CSV/TSV/PDF/DOCX/XLSX/PPTX."
            )

        for enc in ("utf-8", "utf-8-sig", "gb18030", "gbk", "utf-16", "utf-16-le", "utf-16-be"):
            try:
                return raw.decode(enc), f"text({enc})"
            except UnicodeDecodeError:
                continue

        raise ToolExecutionError(f"Unable to decode file: {path.name}. Please convert to UTF-8 text.")

    def _looks_binary(self, raw: bytes) -> bool:
        if not raw:
            return False
        if b"\x00" in raw:
            return True
        sample = raw[:4096]
        ctrl = sum(1 for b in sample if b < 9 or (13 < b < 32))
        return (ctrl / max(1, len(sample))) > 0.12

    def _read_json(self, path: Path) -> str:
        for enc in ("utf-8", "utf-8-sig", "gb18030", "gbk", "utf-16"):
            try:
                text = path.read_text(encoding=enc)
                obj = json.loads(text)
                return json.dumps(obj, ensure_ascii=False, indent=2)
            except UnicodeDecodeError:
                continue
            except json.JSONDecodeError:
                break
        raise ToolExecutionError(f"JSON parse failed: {path.name}")

    def _read_csv_like(self, path: Path) -> str:
        for enc in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
            try:
                return path.read_text(encoding=enc)
            except UnicodeDecodeError:
                continue
        raise ToolExecutionError(f"CSV/TSV encoding not recognized: {path.name}")

    def _read_pdf(self, path: Path) -> str:
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception as exc:
            raise ToolExecutionError("PDF reader dependency missing. Install `pypdf`.") from exc

        try:
            reader = PdfReader(str(path))
            parts: list[str] = []
            for page in reader.pages:
                t = (page.extract_text() or "").strip()
                if t:
                    parts.append(t)
            return "\n\n".join(parts).strip()
        except Exception as exc:
            raise ToolExecutionError(f"PDF parse failed: {path.name}") from exc

    def _read_docx(self, path: Path) -> str:
        try:
            from docx import Document  # type: ignore

            doc = Document(str(path))
            lines = [p.text.strip() for p in doc.paragraphs if (p.text or "").strip()]
            if lines:
                return "\n".join(lines)
        except Exception:
            pass

        try:
            with zipfile.ZipFile(path) as zf:
                xml_bytes = zf.read("word/document.xml")
        except Exception as exc:
            raise ToolExecutionError(f"DOCX read failed: {path.name}") from exc

        try:
            root = ET.fromstring(xml_bytes)
        except ET.ParseError as exc:
            raise ToolExecutionError(f"DOCX parse failed: {path.name}") from exc

        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        lines: list[str] = []
        for para in root.findall(".//w:p", ns):
            texts = [t.text or "" for t in para.findall(".//w:t", ns)]
            line = "".join(texts).strip()
            if line:
                lines.append(line)
        return "\n".join(lines)

    def _read_xlsx(self, path: Path) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore

            wb = load_workbook(filename=str(path), read_only=True, data_only=True)
            out: list[str] = []
            for ws in wb.worksheets:
                out.append(f"[{ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    values = ["" if v is None else str(v) for v in row]
                    if any(v.strip() for v in values):
                        out.append("\t".join(values))
            content = "\n".join(out).strip()
            if content:
                return content
        except Exception:
            pass

        try:
            with zipfile.ZipFile(path) as zf:
                shared = self._xlsx_shared_strings(zf)
                sheets = sorted(
                    n
                    for n in zf.namelist()
                    if n.startswith("xl/worksheets/sheet") and n.endswith(".xml")
                )
                out: list[str] = []
                for sheet in sheets:
                    out.append(f"[{Path(sheet).name}]")
                    out.extend(self._xlsx_sheet_rows(zf.read(sheet), shared))
                return "\n".join(out).strip()
        except Exception as exc:
            raise ToolExecutionError(f"XLSX parse failed: {path.name}") from exc

    def _xlsx_shared_strings(self, zf: zipfile.ZipFile) -> list[str]:
        if "xl/sharedStrings.xml" not in zf.namelist():
            return []
        root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
        ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        vals: list[str] = []
        for si in root.findall(".//s:si", ns):
            vals.append("".join((t.text or "") for t in si.findall(".//s:t", ns)))
        return vals

    def _xlsx_sheet_rows(self, xml_bytes: bytes, shared: list[str]) -> list[str]:
        root = ET.fromstring(xml_bytes)
        ns = {"s": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        rows: list[str] = []
        for row in root.findall(".//s:row", ns):
            vals: list[str] = []
            for cell in row.findall("s:c", ns):
                cell_type = cell.attrib.get("t")
                v = cell.find("s:v", ns)
                if v is None or v.text is None:
                    vals.append("")
                    continue
                txt = v.text
                if cell_type == "s":
                    try:
                        txt = shared[int(txt)]
                    except Exception:
                        pass
                vals.append(txt)
            if any(v.strip() for v in vals):
                rows.append("\t".join(vals))
        return rows

    def _read_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(str(path))
            out: list[str] = []
            for idx, slide in enumerate(prs.slides, start=1):
                out.append(f"[slide{idx}]")
                for shape in slide.shapes:
                    text = (getattr(shape, "text", "") or "").strip()
                    if text:
                        out.append(text)
            content = "\n".join(out).strip()
            if content:
                return content
        except Exception:
            pass

        try:
            with zipfile.ZipFile(path) as zf:
                slides = sorted(
                    n
                    for n in zf.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                )
                out: list[str] = []
                for slide in slides:
                    out.append(f"[{Path(slide).name}]")
                    out.extend(self._pptx_slide_lines(zf.read(slide)))
                return "\n".join(out).strip()
        except Exception as exc:
            raise ToolExecutionError(f"PPTX parse failed: {path.name}") from exc

    def _pptx_slide_lines(self, xml_bytes: bytes) -> list[str]:
        root = ET.fromstring(xml_bytes)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        lines: list[str] = []
        for p in root.findall(".//a:p", ns):
            line = "".join((t.text or "") for t in p.findall(".//a:t", ns)).strip()
            if line:
                lines.append(line)
        return lines

    def _write_text_file(self, args: dict[str, Any]) -> str:
        rel_path = args["path"]
        content = args["content"]
        mode = args.get("mode", "overwrite")

        path = safe_resolve_path(self.root, rel_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        if mode == "append":
            with path.open("a", encoding="utf-8") as f:
                f.write(content)
        else:
            path.write_text(content, encoding="utf-8")

        return json.dumps({"path": str(path), "mode": mode, "written_chars": len(content)}, ensure_ascii=False)

    def _run_shell_command(self, args: dict[str, Any]) -> str:
        cmd = args.get("command")
        if not isinstance(cmd, list) or not cmd or not all(isinstance(x, str) and x.strip() for x in cmd):
            raise ToolExecutionError("command must be a non-empty string array")

        timeout_seconds = int(args.get("timeout_seconds", 8))
        if timeout_seconds < 1 or timeout_seconds > 20:
            raise ToolExecutionError("timeout_seconds must be in [1, 20]")

        command = [x.strip() for x in cmd]
        program = command[0].lower()
        blocked_chars = {"|", "&&", "||", ";", ">", "<", "$(", "`"}
        for token in command:
            if any(sym in token for sym in blocked_chars):
                raise ToolExecutionError("shell meta characters are not allowed")

        allowed = {
            "python": {"--version", "-V", "-m"},
            "py": {"--version", "-V"},
            "where": None,
            "whoami": None,
            "dir": None,
        }
        if program not in allowed:
            raise ToolExecutionError(f"command not allowed: {program}")

        if program in {"python", "py"}:
            if len(command) == 1:
                raise ToolExecutionError("python command requires explicit safe args")
            if command[1] not in allowed[program]:
                raise ToolExecutionError(f"python arg not allowed: {command[1]}")
            if command[1:3] == ["-m", "pip"]:
                if len(command) < 4 or command[3] not in {"show"}:
                    raise ToolExecutionError("only 'python -m pip show <pkg>' is allowed")
                if len(command) != 5:
                    raise ToolExecutionError("pip show requires exactly one package name")
                pkg = command[4]
                if not re.fullmatch(r"[A-Za-z0-9_.-]{1,80}", pkg):
                    raise ToolExecutionError("invalid package name for pip show")
            elif len(command) > 2:
                raise ToolExecutionError("only python --version / -V or python -m pip show <pkg> are allowed")

        if program == "where":
            if len(command) != 2:
                raise ToolExecutionError("where requires exactly one executable name")
            if not re.fullmatch(r"[A-Za-z0-9_.-]{1,80}", command[1]):
                raise ToolExecutionError("invalid target for where")

        if program in {"whoami", "dir"} and len(command) != 1:
            raise ToolExecutionError(f"{program} does not accept extra args in sandbox mode")

        if program == "dir":
            # Use PowerShell Get-ChildItem instead of cmd built-in.
            ps_cmd = ["powershell", "-NoProfile", "-Command", "Get-ChildItem -Name"]
            try:
                completed = subprocess.run(
                    ps_cmd,
                    cwd=str(self.root),
                    capture_output=True,
                    text=True,
                    timeout=timeout_seconds,
                    shell=False,
                )
            except subprocess.TimeoutExpired as exc:
                raise ToolExecutionError(f"command timeout after {timeout_seconds}s") from exc
        else:
            try:
                completed = subprocess.run(
                    command,
                    cwd=str(self.root),
                    capture_output=True,
                    text=True,
                    timeout=timeout_seconds,
                    shell=False,
                )
            except subprocess.TimeoutExpired as exc:
                raise ToolExecutionError(f"command timeout after {timeout_seconds}s") from exc

        return json.dumps(
            {
                "ok": completed.returncode == 0,
                "returncode": completed.returncode,
                "stdout": completed.stdout[:6000],
                "stderr": completed.stderr[:3000],
            },
            ensure_ascii=False,
        )

    def _delegate_subagent(self, args: dict[str, Any]) -> str:
        from llm_client import build_client
        from config import MODEL_NAME

        task = str(args.get("task", "")).strip()
        if not task:
            raise ToolExecutionError("task cannot be empty")
        files = args.get("files", [])
        if files is None:
            files = []
        if not isinstance(files, list):
            raise ToolExecutionError("files must be a string array")

        snippets: list[str] = []
        for rel in files[:5]:
            if not isinstance(rel, str) or not rel.strip():
                continue
            try:
                raw = self._read_text_file({"path": rel.strip(), "max_chars": 2500})
                payload = json.loads(raw)
                snippets.append(
                    f"[{rel.strip()} | {payload.get('detected_format','text')}]\n{payload.get('content','')}"
                )
            except Exception as exc:  # noqa: BLE001
                snippets.append(f"[{rel.strip()}] read failed: {exc}")

        system_prompt = (
            "You are a delegated sub-agent. Solve only the requested subtask clearly and briefly. "
            "If file snippets are provided, prioritize them."
        )
        user_prompt = task
        if snippets:
            user_prompt += "\n\nFile snippets:\n" + "\n\n---\n\n".join(snippets)

        client = build_client()
        resp = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.2,
        )
        answer = (resp.choices[0].message.content or "").strip()
        return json.dumps(
            {
                "ok": True,
                "task": task,
                "used_files": files[:5],
                "answer": answer,
            },
            ensure_ascii=False,
        )

    def _get_current_time(self, args: dict[str, Any]) -> str:
        utc_offset = str(args.get("utc_offset", "") or "").strip()
        if utc_offset:
            tz = self._parse_utc_offset(utc_offset)
            now = datetime.now(tz)
            return json.dumps(
                {
                    "timezone_mode": "utc_offset",
                    "utc_offset": utc_offset,
                    "iso": now.isoformat(timespec="seconds"),
                    "date": now.strftime("%Y-%m-%d"),
                    "time": now.strftime("%H:%M:%S"),
                    "weekday": now.strftime("%A"),
                },
                ensure_ascii=False,
            )

        now = datetime.now().astimezone()
        offset = now.strftime("%z")
        normalized_offset = f"{offset[:3]}:{offset[3:]}" if offset else ""
        return json.dumps(
            {
                "timezone_mode": "local",
                "utc_offset": normalized_offset,
                "iso": now.isoformat(timespec="seconds"),
                "date": now.strftime("%Y-%m-%d"),
                "time": now.strftime("%H:%M:%S"),
                "weekday": now.strftime("%A"),
            },
            ensure_ascii=False,
        )

    def _parse_utc_offset(self, value: str) -> timezone:
        match = re.fullmatch(r"([+-])(\d{2}):(\d{2})", value)
        if not match:
            raise ToolExecutionError("utc_offset must be like +08:00 or -05:00")
        sign, hours_str, minutes_str = match.groups()
        hours = int(hours_str)
        minutes = int(minutes_str)
        if hours > 23 or minutes > 59:
            raise ToolExecutionError("utc_offset out of valid range")
        total_minutes = hours * 60 + minutes
        if sign == "-":
            total_minutes = -total_minutes
        return timezone(timedelta(minutes=total_minutes))

    def _search_web(self, args: dict[str, Any]) -> str:
        query = str(args["query"]).strip()
        max_results = int(args.get("max_results", 5))
        if not query:
            raise ToolExecutionError("search query cannot be empty")

        search_360_key = (
            os.getenv("SEARCH360_API_KEY")
            or os.getenv("QIHOO360_API_KEY")
            or os.getenv("AI360_API_KEY")
        )
        if search_360_key:
            return self._search_with_360(query, max_results, search_360_key)

        serper_key = os.getenv("SERPER_API_KEY")
        if serper_key:
            return self._search_with_serper(query, max_results, serper_key)

        return self._search_with_duckduckgo(query, max_results)

    def _search_with_360(self, query: str, max_results: int, api_key: str) -> str:
        payload = json.dumps(
            {
                "model": "360gpt-pro",
                "messages": [{"role": "user", "content": query}],
                "max_refer_search_items": max(1, min(max_results, 10)),
                "enable_corner_markers": True,
                "enable_web_page_safety": True,
                "stream": False,
            }
        ).encode("utf-8")
        request = urllib.request.Request(
            "https://api.360.cn/v1/search/aisearch",
            data=payload,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            method="POST",
        )
        data = self._load_json_response(request, "360 search")

        answer = self._extract_360_answer_text(data)
        results = self._extract_360_references(data, max_results)
        if not answer and not results:
            raise ToolExecutionError("360 search returned no usable content")

        return json.dumps(
            {
                "query": query,
                "provider": "360_aisearch",
                "answer": answer,
                "results": results,
                "raw": data,
            },
            ensure_ascii=False,
        )

    def _extract_360_answer_text(self, data: dict[str, Any]) -> str:
        candidates: list[str] = []

        if isinstance(data.get("answer"), str):
            candidates.append(str(data["answer"]))
        if isinstance(data.get("content"), str):
            candidates.append(str(data["content"]))
        if isinstance(data.get("output_text"), str):
            candidates.append(str(data["output_text"]))

        choices = data.get("choices")
        if isinstance(choices, list):
            for choice in choices:
                if not isinstance(choice, dict):
                    continue
                message = choice.get("message")
                if isinstance(message, dict):
                    content = message.get("content")
                    if isinstance(content, str):
                        candidates.append(content)
                    elif isinstance(content, list):
                        for item in content:
                            if isinstance(item, dict) and isinstance(item.get("text"), str):
                                candidates.append(item["text"])

        output = data.get("output")
        if isinstance(output, list):
            for item in output:
                if isinstance(item, dict):
                    content = item.get("content")
                    if isinstance(content, list):
                        for part in content:
                            if isinstance(part, dict) and isinstance(part.get("text"), str):
                                candidates.append(part["text"])

        for text in candidates:
            text = text.strip()
            if text:
                return text
        return ""

    def _extract_360_references(self, data: dict[str, Any], max_results: int) -> list[dict[str, str]]:
        refs: list[dict[str, str]] = []
        possible_lists = [
            data.get("references"),
            data.get("refer_search_items"),
            data.get("search_results"),
            data.get("results"),
            data.get("items"),
        ]

        for items in possible_lists:
            if not isinstance(items, list):
                continue
            for item in items:
                if len(refs) >= max_results:
                    return refs
                if not isinstance(item, dict):
                    continue
                title = str(item.get("title") or item.get("name") or "").strip()
                link = str(item.get("url") or item.get("link") or "").strip()
                snippet = str(
                    item.get("summary_ai")
                    or item.get("summary")
                    or item.get("snippet")
                    or item.get("text")
                    or ""
                ).strip()
                if title or link or snippet:
                    refs.append({"title": title, "link": link, "snippet": snippet})

        return refs

    def _search_with_serper(self, query: str, max_results: int, api_key: str) -> str:
        payload = json.dumps({"q": query, "num": max_results}).encode("utf-8")
        request = urllib.request.Request(
            "https://google.serper.dev/search",
            data=payload,
            headers={
                "X-API-KEY": api_key,
                "Content-Type": "application/json",
            },
            method="POST",
        )
        data = self._load_json_response(request, "Serper")

        organic = data.get("organic", [])[:max_results]
        results = [
            {
                "title": item.get("title", ""),
                "link": item.get("link", ""),
                "snippet": item.get("snippet", ""),
            }
            for item in organic
        ]
        return json.dumps(
            {
                "query": query,
                "provider": "serper",
                "answer_box": data.get("answerBox"),
                "knowledge_graph": data.get("knowledgeGraph"),
                "results": results,
            },
            ensure_ascii=False,
        )

    def _search_with_duckduckgo(self, query: str, max_results: int) -> str:
        url = "https://api.duckduckgo.com/?" + urllib.parse.urlencode(
            {
                "q": query,
                "format": "json",
                "no_html": "1",
                "skip_disambig": "0",
            }
        )
        request = urllib.request.Request(
            url,
            headers={"User-Agent": "Mozilla/5.0"},
            method="GET",
        )
        data = self._load_json_response(request, "DuckDuckGo")

        results: list[dict[str, str]] = []
        abstract = str(data.get("AbstractText", "")).strip()
        abstract_url = str(data.get("AbstractURL", "")).strip()
        heading = str(data.get("Heading", "")).strip()
        if abstract:
            results.append(
                {
                    "title": heading or query,
                    "link": abstract_url,
                    "snippet": abstract,
                }
            )

        def collect_topics(items: list[Any]) -> None:
            for item in items:
                if len(results) >= max_results:
                    return
                if isinstance(item, dict) and "Topics" in item:
                    collect_topics(item.get("Topics", []))
                    continue
                if not isinstance(item, dict):
                    continue
                text = str(item.get("Text", "")).strip()
                link = str(item.get("FirstURL", "")).strip()
                if text:
                    title = text.split(" - ", 1)[0]
                    results.append({"title": title, "link": link, "snippet": text})

        collect_topics(data.get("RelatedTopics", []))

        if not results:
            raise ToolExecutionError("No usable search results found")

        return json.dumps(
            {
                "query": query,
                "provider": "duckduckgo_instant_answer",
                "results": results[:max_results],
            },
            ensure_ascii=False,
        )

    def _load_json_response(self, request: urllib.request.Request, provider_label: str) -> dict[str, Any]:
        try:
            with urllib.request.urlopen(request, timeout=20) as response:
                return json.loads(response.read().decode("utf-8"))
        except urllib.error.HTTPError as exc:
            body = ""
            try:
                body = exc.read().decode("utf-8", errors="ignore")
            except Exception:
                pass
            raise ToolExecutionError(f"{provider_label} HTTP {exc.code}: {body[:200]}") from exc
        except urllib.error.URLError as exc:
            raise ToolExecutionError(f"{provider_label} network error: {exc.reason}") from exc
        except json.JSONDecodeError as exc:
            raise ToolExecutionError(f"{provider_label} returned invalid JSON: {exc}") from exc

    def _get_weather(self, args: dict[str, Any]) -> str:
        location = str(args["location"]).strip()
        if not location:
            raise ToolExecutionError("location cannot be empty")

        geocode_url = (
            "https://geocoding-api.open-meteo.com/v1/search?"
            + urllib.parse.urlencode({"name": location, "count": 1, "language": "zh", "format": "json"})
        )
        geocode_req = urllib.request.Request(geocode_url, method="GET", headers={"User-Agent": "Mozilla/5.0"})
        geo = self._load_json_response(geocode_req, "open-meteo geocoding")
        results = geo.get("results")
        if not isinstance(results, list) or not results:
            raise ToolExecutionError(f"Location not found: {location}")

        place = results[0]
        lat = place.get("latitude")
        lon = place.get("longitude")
        if lat is None or lon is None:
            raise ToolExecutionError(f"Location missing coordinates: {location}")

        weather_params = {
            "latitude": lat,
            "longitude": lon,
            "current": "temperature_2m,relative_humidity_2m,apparent_temperature,precipitation,weather_code,wind_speed_10m",
            "daily": "weather_code,temperature_2m_max,temperature_2m_min,precipitation_sum",
            "timezone": "auto",
            "forecast_days": 3,
        }
        weather_url = "https://api.open-meteo.com/v1/forecast?" + urllib.parse.urlencode(weather_params)
        weather_req = urllib.request.Request(weather_url, method="GET", headers={"User-Agent": "Mozilla/5.0"})
        data = self._load_json_response(weather_req, "open-meteo weather")

        current = data.get("current", {})
        daily = data.get("daily", {})
        dates = daily.get("time", []) or []
        max_t = daily.get("temperature_2m_max", []) or []
        min_t = daily.get("temperature_2m_min", []) or []
        rain = daily.get("precipitation_sum", []) or []
        d_code = daily.get("weather_code", []) or []

        forecast = []
        for i in range(min(3, len(dates))):
            forecast.append(
                {
                    "date": dates[i],
                    "weather": self._weather_code_to_text(d_code[i] if i < len(d_code) else None),
                    "temp_max_c": max_t[i] if i < len(max_t) else None,
                    "temp_min_c": min_t[i] if i < len(min_t) else None,
                    "precipitation_mm": rain[i] if i < len(rain) else None,
                }
            )

        return json.dumps(
            {
                "query_location": location,
                "resolved_location": {
                    "name": place.get("name"),
                    "country": place.get("country"),
                    "admin1": place.get("admin1"),
                    "latitude": lat,
                    "longitude": lon,
                },
                "current": {
                    "temperature_c": current.get("temperature_2m"),
                    "feels_like_c": current.get("apparent_temperature"),
                    "humidity": current.get("relative_humidity_2m"),
                    "wind_speed_kmh": current.get("wind_speed_10m"),
                    "precipitation_mm": current.get("precipitation"),
                    "weather": self._weather_code_to_text(current.get("weather_code")),
                },
                "forecast": forecast,
            },
            ensure_ascii=False,
        )

    def _weather_code_to_text(self, code: Any) -> str:
        mapping = {
            0: "晴朗",
            1: "少云",
            2: "多云",
            3: "阴天",
            45: "雾",
            48: "冻雾",
            51: "小毛毛雨",
            53: "毛毛雨",
            55: "强毛毛雨",
            61: "小雨",
            63: "中雨",
            65: "大雨",
            71: "小雪",
            73: "中雪",
            75: "大雪",
            80: "阵雨",
            81: "较强阵雨",
            82: "强阵雨",
            95: "雷暴",
        }
        if isinstance(code, (int, float)):
            return mapping.get(int(code), f"未知天气代码({int(code)})")
        return "未知"
