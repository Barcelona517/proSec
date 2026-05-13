from pptx import Presentation
prs = Presentation("Cppt/C语言特性深度解析.pptx")
print(f"总页数: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\n--- 第{i+1}页 ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text[:200])
