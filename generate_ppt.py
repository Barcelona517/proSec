from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT_COLOR = RGBColor(0x00, 0xAD, 0xB5)
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RGBColor(0xCC, 0xCC, 0xCC)

def add_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_box(slide, text, left=0.5, top=0.3, width=12, height=1.2, font_size=36, color=TITLE_COLOR):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return txBox

def add_bullet_box(slide, items, left=0.8, top=1.8, width=11.5, height=5, font_size=20, color=TEXT_COLOR):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.level = 0
    return txBox

# ===== Slide 1: Title =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)
add_title_box(slide1, "C 语言核心特性", left=1, top=2, width=11, height=1.5, font_size=48)
add_title_box(slide1, "高效 · 灵活 · 底层控制", left=1, top=3.5, width=11, height=1, font_size=28, color=ACCENT_COLOR)
add_title_box(slide1, "2025年4月", left=1, top=4.8, width=11, height=0.8, font_size=18, color=TEXT_COLOR)

# ===== Slide 2: 过程式编程 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title_box(slide2, "特性一：过程式编程范式", font_size=40)
items2 = [
    "• C 语言是典型的过程式（结构化）编程语言",
    "• 程序由函数（Function）组成，main() 是入口",
    "• 支持顺序、选择（if/switch）、循环（for/while/do-while）三种基本结构",
    "• 采用自顶向下、逐步细化的设计方法",
    "• 代码可读性强，逻辑清晰，适合系统级开发",
    "",
    "示例：",
    '  int main() {',
    '      int a = 10, b = 20;',
    '      printf("Sum = %d\\n", a + b);',
    '      return 0;',
    '  }'
]
add_bullet_box(slide2, items2, font_size=18)

# ===== Slide 3: 指针与内存管理 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title_box(slide3, "特性二：指针与直接内存管理", font_size=40)
items3 = [
    "• 指针（Pointer）是 C 语言的灵魂，直接操作内存地址",
    "• 支持指针运算：++、--、加减偏移量",
    "• 动态内存分配：malloc()、calloc()、realloc()、free()",
    "• 指针与数组紧密关联，数组名即首地址",
    "• 函数指针可实现回调机制",
    "",
    "示例：",
    '  int x = 42;',
    '  int *p = &x;',
    '  printf("%d", *p);   // 输出 42',
    '',
    "⚠️ 需手动管理内存，避免内存泄漏和野指针"
]
add_bullet_box(slide3, items3, font_size=18)

# ===== Slide 4: 底层硬件操作 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title_box(slide4, "特性三：底层硬件操作能力", font_size=40)
items4 = [
    "• 位运算：&（与）、|（或）、^（异或）、~（取反）、<<（左移）、>>（右移）",
    "• 直接操作寄存器、端口、内存映射 I/O",
    "• 联合体（union）实现同一内存区域多种解释",
    "• 位域（Bit-field）精确控制结构体成员的比特位数",
    "• 内联汇编（asm）可在 C 代码中嵌入汇编指令",
    "",
    "应用场景：",
    "  → 嵌入式系统、操作系统内核、驱动程序开发",
    "  → 单片机编程、实时控制系统"
]
add_bullet_box(slide4, items4, font_size=18)

# ===== Slide 5: 总结 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title_box(slide5, "总结：C 语言的独特优势", font_size=40)
items5 = [
    "✅ 高效性：编译后代码接近汇编，执行效率极高",
    "✅ 可移植性：几乎支持所有平台和架构",
    "✅ 灵活性：指针、位运算、宏定义等赋予极大自由度",
    "✅ 影响力：C 语言深刻影响了 C++、Java、C#、Python 等语言",
    "",
    "📌 适合领域：",
    "   • 操作系统（Linux、Windows 内核）",
    "   • 嵌入式系统与物联网设备",
    "   • 高性能计算与数据库引擎",
    "   • 编译器和解释器开发",
    "",
    "\"C 语言——让你离机器更近，离抽象更远。\""
]
add_bullet_box(slide5, items5, font_size=18)

# Save
output_path = "C语言核心特性.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
